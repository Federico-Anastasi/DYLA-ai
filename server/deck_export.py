"""Exports deck.json in two formats: pptx (python-pptx) and a self-contained HTML page.

deck.json is the source document (see schemas/deck.schema.json): a project slide deck
(kick-off, status review, estimate walkthrough, demo) meant to be regenerated
periodically (see the 'auto' flag on each slide) without losing the hand-written parts.
This module only renders whatever it finds in the JSON — the incremental regeneration
logic (which slides are 'auto', how they get refreshed from estimate.json/timeline.json)
lives in the /deck skill, not here.

Module style mirrors server/mockup_export.py:
- one `_render_*`/`_slide_*` function per layout;
- CSS/JS inlined in the HTML export, zero external resources (no CDN, no remote fonts,
  images embedded as data URIs);
- color themes defined in a small table at module level.

Rendering choices:
- 16:9 in both formats (pptx: slide dimensions set explicitly; html: the ratio is
  cosmetic, content scrolls when it has to).
- 'cover' and 'section' fill the whole slide with the theme background; every other
  layout gets a colored title bar at the top and the content below (same composition
  in both exports).
- 'text' supports a minimal markdown: paragraphs separated by a blank line, lines
  starting with '- '/'* ' become bullets, **bold**/*italic* become formatted runs/tags —
  and nothing else (no nested headings, no links, no inline tables), per the schema.
- 'timeline': past 6 milestones, laying them out on a single horizontal line becomes
  unreadable (labels overlap), so it degrades into a Name/Date/Status table — same
  threshold in pptx and in html.
- a missing 'image' never breaks the export: a placeholder shows the path that was tried.
"""
from __future__ import annotations

import base64
import html
import io
import mimetypes
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .config import PROJECTS_DIR
from .documents import load_doc
from .exports import DocNotFound

__all__ = ["deck_pptx", "deck_html"]


def _require_deck(project: str) -> dict:
    data = load_doc(project, "deck")
    if data is None:
        raise DocNotFound(f"deck.json not found for project '{project}'")
    return data


def _project_image(meta: dict, rel: str | None) -> Path | None:
    """A slide's image, only if it really is inside that project's folder.

    `slides[].image` is a path written by the model into a JSON file, and the schema puts
    no constraint on it. Joined naively, `../../server/config.yaml` walks out of the
    project — and on Windows an absolute path makes pathlib discard the prefix entirely,
    so the join silently becomes "anywhere on this disk". The HTML export then embeds
    whatever it found as base64 inside a file made to be sent to a client.

    Nobody is attacking a local app on purpose. But this is a file-read primitive driven
    by generated content, and it costs three lines to close.
    """
    if not rel:
        return None
    root = (PROJECTS_DIR / meta["project"]).resolve()
    try:
        path = (root / rel).resolve()
    except (OSError, ValueError):
        return None
    return path if path.is_relative_to(root) and path.is_file() else None


def _esc(s) -> str:
    return html.escape(str(s if s is not None else ""), quote=True)


# ── color theme, same three themes as the mockup library ──────────────────

_THEMES = {"standard": "2F5D8C", "compact": "1F7A6C", "plain": "5B5470"}
_DEFAULT_PRIMARY = "3B5166"  # neutral blue-grey for an unknown theme name


def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _rgb_to_hex(rgb: tuple[int, int, int]) -> str:
    return "".join(f"{max(0, min(255, c)):02X}" for c in rgb)


def _lighten(h: str, factor: float) -> str:
    r, g, b = _hex_to_rgb(h)
    return _rgb_to_hex((
        int(r + (255 - r) * factor),
        int(g + (255 - g) * factor),
        int(b + (255 - b) * factor),
    ))


def _darken(h: str, factor: float) -> str:
    r, g, b = _hex_to_rgb(h)
    return _rgb_to_hex((int(r * (1 - factor)), int(g * (1 - factor)), int(b * (1 - factor))))


def _theme_of(theme: str) -> dict:
    primary = _THEMES.get(theme, _DEFAULT_PRIMARY)
    return {"primary": primary, "light": _lighten(primary, 0.88), "dark": _darken(primary, 0.30)}


# Milestone status -> color (see schemas/deck.schema.json, slides[].milestones[].status).
# 'in_progress' has no fixed color: it picks up the deck theme (see _status_hex).
_STATUS_HEX = {"planned": "9CA3AF", "completed": "22C55E", "at_risk": "DC2626"}
_STATUS_LABEL = {"planned": "Planned", "in_progress": "In progress",
                  "completed": "Completed", "at_risk": "At risk"}


def _status_hex(status: str, theme: dict) -> str:
    if status == "in_progress":
        return theme["primary"]
    return _STATUS_HEX.get(status, _STATUS_HEX["planned"])


# Past this threshold the horizontal timeline degrades into a table (see module docstring).
_TIMELINE_MAX_INLINE = 6


# ── minimal markdown shared by pptx and html ('text' layout and bullets) ───

_MD_SPAN = re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*)")
_MD_BOLD = re.compile(r"\*\*([^*]+)\*\*")
_MD_ITALIC = re.compile(r"\*([^*]+)\*")


def _parse_simple_markdown(text: str) -> list[tuple[str, bool]]:
    """Splits 'text' into (content, is_bullet) blocks. Paragraphs separated by a blank
    line are joined back onto a single line; lines starting with '- '/'* ' are single
    bullets."""
    out: list[tuple[str, bool]] = []
    buffer: list[str] = []

    def flush():
        if buffer:
            out.append((" ".join(buffer).strip(), False))
            buffer.clear()

    for line in (text or "").strip().splitlines():
        stripped = line.strip()
        if not stripped:
            flush()
            continue
        if stripped.startswith("- ") or stripped.startswith("* "):
            flush()
            out.append((stripped[2:].strip(), True))
        else:
            buffer.append(stripped)
    flush()
    return out


def _md_inline_html(text: str) -> str:
    escaped = _esc(text)
    escaped = _MD_BOLD.sub(r"<strong>\1</strong>", escaped)
    escaped = _MD_ITALIC.sub(r"<em>\1</em>", escaped)
    return escaped


# ═══════════════════════════════════════════════════════════════════════════
# PPTX
# ═══════════════════════════════════════════════════════════════════════════

SLIDE_W = Emu(12192000)  # 13.333 in — 16:9
SLIDE_H = Emu(6858000)   # 7.5 in


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 'Blank'


def _textbox(slide, left, top, width, height, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return box, tf


def _set_run(run, text, *, size=18, color="1A1A1A", bold=False, italic=False, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = RGBColor.from_string(color)


def _line(tf, text, *, first=False, align=PP_ALIGN.LEFT, **run_kwargs):
    """Single-line paragraph, no markdown interpretation (titles, footers, cells)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    _set_run(p.add_run(), text, **run_kwargs)
    return p


def _markdown_line(tf, text, *, first=False, align=PP_ALIGN.LEFT, prefix=None, **run_kwargs):
    """Paragraph where **bold**/*italic* become separate runs; 'prefix' (a bullet, say)
    is a run of its own and is never treated as markdown."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if prefix:
        _set_run(p.add_run(), prefix, **{**run_kwargs, "bold": True})
    for tok in _MD_SPAN.split(text):
        if not tok:
            continue
        if tok.startswith("**") and tok.endswith("**") and len(tok) >= 4:
            _set_run(p.add_run(), tok[2:-2], **{**run_kwargs, "bold": True})
        elif tok.startswith("*") and tok.endswith("*") and len(tok) >= 2:
            _set_run(p.add_run(), tok[1:-1], **{**run_kwargs, "italic": True})
        else:
            _set_run(p.add_run(), tok, **run_kwargs)
    return p


def _content_header(slide, theme: dict, title: str):
    """Colored title bar at the top, shared by every layout except 'cover'/'section'."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(theme["primary"])
    bar.line.fill.background()
    box, tf = _textbox(slide, Inches(0.6), Inches(0.1), SLIDE_W - Inches(1.2), Inches(0.9),
                        anchor=MSO_ANCHOR.MIDDLE)
    _line(tf, title, first=True, size=28, bold=True, color="FFFFFF")


def _slide_cover(prs, s, meta, theme):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(theme["primary"])
    bg.line.fill.background()
    _, tf = _textbox(slide, Inches(1), Inches(2.6), SLIDE_W - Inches(2), Inches(1.4),
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _line(tf, s["title"], first=True, size=40, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    if s.get("subtitle"):
        _, tf2 = _textbox(slide, Inches(1), Inches(4.05), SLIDE_W - Inches(2), Inches(0.8),
                           align=PP_ALIGN.CENTER)
        _line(tf2, s["subtitle"], first=True, size=20, color="FFFFFF", align=PP_ALIGN.CENTER)
    footer = f'{meta.get("client", "")} — {meta.get("date", "")}'
    _, tf3 = _textbox(slide, Inches(1), SLIDE_H - Inches(1.0), SLIDE_W - Inches(2), Inches(0.5),
                       align=PP_ALIGN.CENTER)
    _line(tf3, footer, first=True, size=14, color="FFFFFF", align=PP_ALIGN.CENTER)
    return slide


def _slide_section(prs, s, meta, theme):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(theme["dark"])
    bg.line.fill.background()
    _, tf = _textbox(slide, Inches(1), Inches(3.0), SLIDE_W - Inches(2), Inches(1.2),
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _line(tf, s["title"], first=True, size=34, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    if s.get("subtitle"):
        _, tf2 = _textbox(slide, Inches(1), Inches(4.15), SLIDE_W - Inches(2), Inches(0.7),
                           align=PP_ALIGN.CENTER)
        _line(tf2, s["subtitle"], first=True, size=18, color="FFFFFF", align=PP_ALIGN.CENTER)
    return slide


def _slide_list(prs, s, meta, theme):
    slide = _blank_slide(prs)
    _content_header(slide, theme, s["title"])
    _, tf = _textbox(slide, Inches(0.8), Inches(1.5), SLIDE_W - Inches(1.6), SLIDE_H - Inches(2.0))
    for i, b in enumerate(s.get("bullets") or []):
        p = _markdown_line(tf, b, first=(i == 0), prefix="•  ", size=20, color="1A1A1A")
        p.space_after = Pt(14)
    return slide


def _slide_table(prs, s, meta, theme):
    slide = _blank_slide(prs)
    _content_header(slide, theme, s["title"])
    # "table" is not required by the schema (only id/layout/title are), unlike every other
    # layout's own field, which is always read with .get(...) or a fallback — a slide typed
    # "table" with nothing in it must not take the whole export down with a KeyError. Zero
    # columns is not a table python-pptx can even build, so an empty one gets a placeholder
    # instead, the same way a missing image does below.
    table_data = s.get("table") or {}
    headers, rows = table_data.get("headers") or [], table_data.get("rows") or []
    if not headers:
        _, tf = _textbox(slide, Inches(0.8), Inches(1.5), SLIDE_W - Inches(1.6), Inches(1.0))
        _line(tf, "No table data", first=True, size=16, italic=True, color="666666")
        return slide
    gframe = slide.shapes.add_table(len(rows) + 1, len(headers),
                                     Inches(0.6), Inches(1.5), SLIDE_W - Inches(1.2), SLIDE_H - Inches(2.2))
    table = gframe.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(theme["primary"])
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(14)
                r.font.color.rgb = RGBColor.from_string("FFFFFF")
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
    return slide


def _slide_kpi(prs, s, meta, theme):
    slide = _blank_slide(prs)
    _content_header(slide, theme, s["title"])
    kpis = s.get("kpi") or []
    n = len(kpis) or 1
    margin, gap = Inches(0.6), Inches(0.3)
    card_w = Emu(int((SLIDE_W - 2 * margin - gap * (n - 1)) / n))
    card_h, top = Inches(3.0), Inches(2.4)
    for i, k in enumerate(kpis):
        left = Emu(int(margin + i * (card_w + gap)))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(theme["light"])
        card.line.color.rgb = RGBColor.from_string(theme["primary"])
        card.line.width = Pt(1)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _line(tf, k["value"], first=True, size=32, bold=True, color=theme["dark"], align=PP_ALIGN.CENTER)
        _line(tf, k["label"], size=14, color="333333", align=PP_ALIGN.CENTER)
        if k.get("note"):
            _line(tf, k["note"], size=10, italic=True, color="666666", align=PP_ALIGN.CENTER)
    return slide


def _render_milestones_table_pptx(slide, theme, milestones):
    gframe = slide.shapes.add_table(len(milestones) + 1, 3,
                                     Inches(0.8), Inches(1.6), SLIDE_W - Inches(1.6), SLIDE_H - Inches(2.4))
    table = gframe.table
    for c, h in enumerate(["Milestone", "Date", "Status"]):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(theme["primary"])
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.color.rgb = RGBColor.from_string("FFFFFF")
    for ri, m in enumerate(milestones, start=1):
        vals = [m["name"], m.get("date") or "", _STATUS_LABEL.get(m.get("status"), "")]
        for ci, val in enumerate(vals):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)


def _slide_timeline(prs, s, meta, theme):
    slide = _blank_slide(prs)
    _content_header(slide, theme, s["title"])
    milestones = s.get("milestones") or []
    if len(milestones) > _TIMELINE_MAX_INLINE:
        # Threshold: past 6 milestones the labels on a single horizontal row overlap and
        # the layout stops being readable -> degrade to a Name/Date/Status table (same
        # rule in html).
        _render_milestones_table_pptx(slide, theme, milestones)
        return slide
    if not milestones:
        return slide

    line_y = Inches(3.6)
    left_margin, right_margin = Inches(1.0), Inches(1.0)
    usable = SLIDE_W - left_margin - right_margin
    n = len(milestones)
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left_margin, line_y,
                                            SLIDE_W - right_margin, line_y)
    connector.line.color.rgb = RGBColor.from_string("CBD5E1")
    connector.line.width = Pt(2)

    dot_d = Inches(0.22)
    box_w = Inches(1.8)
    for i, m in enumerate(milestones):
        cx = left_margin + (usable * i / (n - 1) if n > 1 else usable / 2)
        color = _status_hex(m.get("status", "planned"), theme)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - dot_d / 2)), Emu(int(line_y - dot_d / 2)),
                                      dot_d, dot_d)
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor.from_string(color)
        dot.line.color.rgb = RGBColor.from_string("FFFFFF")
        dot.line.width = Pt(1.5)
        # labels alternate above and below the line so they do not collide
        label_top = line_y - Inches(1.05) if i % 2 == 0 else line_y + Inches(0.35)
        _, tf = _textbox(slide, Emu(int(cx - box_w / 2)), label_top, box_w, Inches(0.95),
                          align=PP_ALIGN.CENTER)
        _line(tf, m["name"], first=True, size=13, bold=True, color="1A1A1A", align=PP_ALIGN.CENTER)
        if m.get("date"):
            _line(tf, m["date"], size=11, color="666666", align=PP_ALIGN.CENTER)
    return slide


def _slide_text(prs, s, meta, theme):
    slide = _blank_slide(prs)
    _content_header(slide, theme, s["title"])
    _, tf = _textbox(slide, Inches(0.8), Inches(1.5), SLIDE_W - Inches(1.6), SLIDE_H - Inches(2.0))
    blocks = _parse_simple_markdown(s.get("text") or "")
    for i, (text, is_bullet) in enumerate(blocks):
        p = _markdown_line(tf, text, first=(i == 0), prefix=("•  " if is_bullet else None),
                            size=16, color="1A1A1A")
        p.space_after = Pt(10)
    return slide


def _slide_image(prs, s, meta, theme):
    slide = _blank_slide(prs)
    _content_header(slide, theme, s["title"])
    area_left, area_top = Inches(0.8), Inches(1.5)
    area_w, area_h = SLIDE_W - Inches(1.6), SLIDE_H - Inches(2.0)
    rel = s.get("image")
    img_path = _project_image(meta, rel)
    if img_path:
        from PIL import Image  # already a dependency, pulled in by python-pptx

        with Image.open(img_path) as im:
            iw, ih = im.size
        ratio = min(area_w / iw, area_h / ih)
        w, h = Emu(int(iw * ratio)), Emu(int(ih * ratio))
        left = Emu(int(area_left + (area_w - w) / 2))
        top = Emu(int(area_top + (area_h - h) / 2))
        slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, area_left, area_top, area_w, area_h)
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor.from_string("EEEEEE")
        ph.line.color.rgb = RGBColor.from_string("999999")
        tf = ph.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _line(tf, f'Image not found: {rel or "(no path given)"}',
              first=True, size=16, italic=True, color="666666", align=PP_ALIGN.CENTER)
    return slide


_PPTX_RENDERERS = {
    "cover": _slide_cover,
    "section": _slide_section,
    "list": _slide_list,
    "table": _slide_table,
    "kpi": _slide_kpi,
    "timeline": _slide_timeline,
    "text": _slide_text,
    "image": _slide_image,
}


def deck_pptx(project: str) -> bytes:
    data = _require_deck(project)
    meta = data["meta"]
    theme = _theme_of(meta.get("theme", ""))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for s in data["slides"]:
        slide = _PPTX_RENDERERS[s["layout"]](prs, s, meta, theme)
        note = s.get("speaker_notes")
        if note:
            slide.notes_slide.notes_text_frame.text = note

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════
# HTML
# ═══════════════════════════════════════════════════════════════════════════

def _image_data_uri(path: Path) -> str:
    mime = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
    data = base64.b64encode(path.read_bytes()).decode("ascii")
    return f"data:{mime};base64,{data}"


def _html_cover(s, meta, theme) -> str:
    sub = f'<p class="cover-subtitle">{_esc(s["subtitle"])}</p>' if s.get("subtitle") else ""
    footer = f'{_esc(meta.get("client", ""))} — {_esc(meta.get("date", ""))}'
    return (f'<div class="cover"><h1>{_esc(s["title"])}</h1>{sub}'
            f'<div class="cover-footer">{footer}</div></div>')


def _html_section(s, meta, theme) -> str:
    sub = f'<p class="section-subtitle">{_esc(s["subtitle"])}</p>' if s.get("subtitle") else ""
    return f'<div class="section-divider"><h1>{_esc(s["title"])}</h1>{sub}</div>'


def _html_header(title: str) -> str:
    return f'<div class="slide-header">{_esc(title)}</div>'


def _html_list(s, meta, theme) -> str:
    items = "".join(f"<li>{_md_inline_html(b)}</li>" for b in (s.get("bullets") or []))
    return f'{_html_header(s["title"])}<div class="slide-body"><ul class="bullets">{items}</ul></div>'


def _html_table(s, meta, theme) -> str:
    # See _slide_table's comment: "table" is optional in the schema, unlike every other
    # layout's own field.
    table_data = s.get("table") or {}
    headers, rows_data = table_data.get("headers") or [], table_data.get("rows") or []
    if not headers:
        return (f'{_html_header(s["title"])}<div class="slide-body">'
                f'<div class="image-placeholder">No table data</div></div>')
    thead = "".join(f"<th>{_esc(h)}</th>" for h in headers)
    rows = "".join("<tr>" + "".join(f"<td>{_esc(v)}</td>" for v in row) + "</tr>" for row in rows_data)
    return (f'{_html_header(s["title"])}<div class="slide-body"><div class="table-wrap">'
            f"<table><thead><tr>{thead}</tr></thead><tbody>{rows}</tbody></table></div></div>")


def _html_kpi(s, meta, theme) -> str:
    cards = []
    for k in s.get("kpi") or []:
        note = f'<div class="kpi-note">{_esc(k["note"])}</div>' if k.get("note") else ""
        cards.append(f'<div class="kpi-card"><div class="kpi-value">{_esc(k["value"])}</div>'
                      f'<div class="kpi-label">{_esc(k["label"])}</div>{note}</div>')
    return f'{_html_header(s["title"])}<div class="slide-body"><div class="kpi-row">{"".join(cards)}</div></div>'


def _html_milestones_table(s) -> str:
    rows = "".join(
        f'<tr><td>{_esc(m["name"])}</td><td>{_esc(m.get("date") or "")}</td>'
        f'<td>{_esc(_STATUS_LABEL.get(m.get("status"), ""))}</td></tr>'
        for m in s.get("milestones") or []
    )
    return (f'{_html_header(s["title"])}<div class="slide-body"><div class="table-wrap">'
            f"<table><thead><tr><th>Milestone</th><th>Date</th><th>Status</th></tr></thead>"
            f"<tbody>{rows}</tbody></table></div></div>")


def _html_timeline(s, meta, theme) -> str:
    milestones = s.get("milestones") or []
    if len(milestones) > _TIMELINE_MAX_INLINE:
        return _html_milestones_table(s)
    items = []
    for m in milestones:
        status = m.get("status", "planned")
        date_html = f'<div class="tl-date">{_esc(m["date"])}</div>' if m.get("date") else ""
        items.append(
            f'<div class="tl-item"><div class="tl-dot tl-dot-{_esc(status)}" '
            f'title="{_esc(_STATUS_LABEL.get(status, ""))}"></div>'
            f'<div class="tl-name">{_esc(m["name"])}</div>{date_html}</div>'
        )
    return f'{_html_header(s["title"])}<div class="slide-body"><div class="timeline-row">{"".join(items)}</div></div>'


def _html_text(s, meta, theme) -> str:
    blocks = _parse_simple_markdown(s.get("text") or "")
    parts, in_ul = [], False
    for text, is_bullet in blocks:
        if is_bullet:
            if not in_ul:
                parts.append('<ul class="bullets">')
                in_ul = True
            parts.append(f"<li>{_md_inline_html(text)}</li>")
        else:
            if in_ul:
                parts.append("</ul>")
                in_ul = False
            parts.append(f"<p>{_md_inline_html(text)}</p>")
    if in_ul:
        parts.append("</ul>")
    return f'{_html_header(s["title"])}<div class="slide-body slide-text">{"".join(parts)}</div>'


def _html_image(s, meta, theme) -> str:
    rel = s.get("image")
    img_path = _project_image(meta, rel)
    if img_path:
        body = f'<img class="slide-image" src="{_image_data_uri(img_path)}" alt="{_esc(s["title"])}">'
    else:
        missing = _esc(rel or "(no path given)")
        body = f'<div class="image-placeholder">Image not found: {missing}</div>'
    return f'{_html_header(s["title"])}<div class="slide-body">{body}</div>'


_HTML_RENDERERS = {
    "cover": _html_cover,
    "section": _html_section,
    "list": _html_list,
    "table": _html_table,
    "kpi": _html_kpi,
    "timeline": _html_timeline,
    "text": _html_text,
    "image": _html_image,
}


def _render_slide_html(s: dict, meta: dict, theme: dict, index: int) -> str:
    layout = s["layout"]
    body = _HTML_RENDERERS[layout](s, meta, theme)
    return (f'<section class="slide slide-{_esc(layout)}" data-index="{index}" '
            f'data-title="{_esc(s.get("title", ""))}">{body}</section>')


def _html_css(theme: dict) -> str:
    return f"""
:root {{
  --primary: #{theme["primary"]};
  --primary-light: #{theme["light"]};
  --primary-dark: #{theme["dark"]};
  --dot-completed: #22C55E;
  --dot-at_risk: #DC2626;
  --dot-planned: #9CA3AF;
}}
* {{ box-sizing: border-box; }}
html, body {{ margin: 0; padding: 0; }}
body {{
  font-family: Calibri, "Segoe UI", Arial, sans-serif;
  background: #eef0f2;
  color: #1a1a1a;
}}
@media (prefers-color-scheme: dark) {{
  body {{ background: #14161a; color: #e6e8eb; }}
}}
.deck {{ display: flex; flex-direction: column; gap: 24px; padding: 24px; max-width: 1100px; margin: 0 auto; }}
.slide {{
  background: #ffffff; color: #1a1a1a; border-radius: 10px; box-shadow: 0 1px 4px rgba(0,0,0,.18);
  padding: 0; min-height: 260px; overflow: hidden; aspect-ratio: 16 / 9;
}}
.slide-header {{
  background: var(--primary); color: #fff; font-size: 22px; font-weight: bold;
  padding: 18px 28px;
}}
.slide-body {{ padding: 24px 32px; overflow: auto; }}
.cover, .section-divider {{
  height: 100%; display: flex; flex-direction: column; align-items: center; justify-content: center;
  text-align: center; color: #fff; padding: 40px;
}}
.cover {{ background: var(--primary); }}
.section-divider {{ background: var(--primary-dark); }}
.cover h1 {{ font-size: 40px; margin: 0 0 12px; }}
.section-divider h1 {{ font-size: 32px; margin: 0 0 10px; }}
.cover-subtitle, .section-subtitle {{ font-size: 18px; opacity: .9; margin: 0; }}
.cover-footer {{ position: absolute; bottom: 28px; font-size: 13px; opacity: .85; }}
.bullets {{ margin: 0; padding-left: 22px; line-height: 1.7; font-size: 17px; }}
.table-wrap {{ overflow-x: auto; }}
table {{ width: 100%; border-collapse: collapse; font-size: 14px; }}
th, td {{ border: 1px solid #d8dbe0; padding: 8px 10px; text-align: left; }}
th {{ background: var(--primary); color: #fff; }}
.kpi-row {{ display: flex; flex-wrap: wrap; gap: 16px; }}
.kpi-card {{
  flex: 1 1 160px; background: var(--primary-light); border: 1px solid var(--primary);
  border-radius: 8px; padding: 18px; text-align: center;
}}
.kpi-value {{ font-size: 28px; font-weight: bold; color: var(--primary-dark); }}
.kpi-label {{ font-size: 13px; color: #444; margin-top: 4px; }}
.kpi-note {{ font-size: 11px; color: #777; font-style: italic; margin-top: 6px; }}
.timeline-row {{ display: flex; align-items: flex-start; justify-content: space-between; padding: 40px 10px; position: relative; }}
.timeline-row::before {{
  content: ""; position: absolute; left: 20px; right: 20px; top: 50px; height: 2px; background: #cbd5e1;
}}
.tl-item {{ flex: 1; text-align: center; position: relative; padding: 0 6px; }}
.tl-dot {{ width: 18px; height: 18px; border-radius: 50%; margin: 0 auto 8px; border: 2px solid #fff; box-shadow: 0 0 0 1px #cbd5e1; }}
.tl-dot-planned {{ background: var(--dot-planned); }}
.tl-dot-completed {{ background: var(--dot-completed); }}
.tl-dot-at_risk {{ background: var(--dot-at_risk); }}
.tl-dot-in_progress {{ background: var(--primary); }}
.tl-name {{ font-size: 13px; font-weight: bold; }}
.tl-date {{ font-size: 11px; color: #777; margin-top: 2px; }}
.slide-text p {{ font-size: 16px; line-height: 1.6; }}
.image-placeholder {{
  border: 2px dashed #999; border-radius: 6px; padding: 40px; text-align: center;
  color: #666; font-style: italic;
}}
.slide-image {{ max-width: 100%; max-height: 100%; display: block; margin: 0 auto; }}

/* presentation mode: switched on from JS. By default the deck is a single column, so it
   stays readable and printable even without JS */
body.presenting {{ overflow: hidden; }}
body.presenting .deck {{ display: block; max-width: none; padding: 0; }}
body.presenting .slide {{
  display: none; position: fixed; inset: 0; width: 100vw; height: 100vh;
  border-radius: 0; box-shadow: none; aspect-ratio: auto;
}}
body.presenting .slide.active {{ display: block; }}
body.presenting .slide-body {{ height: calc(100% - 64px); }}

.deck-nav {{
  position: fixed; left: 0; right: 0; bottom: 18px; display: none; align-items: center;
  justify-content: center; gap: 18px; z-index: 20;
}}
body.presenting .deck-nav {{ display: flex; }}
.deck-nav button {{
  background: var(--primary); color: #fff; border: none; border-radius: 50%; width: 40px; height: 40px;
  font-size: 20px; cursor: pointer;
}}
#deck-indicator {{ background: rgba(0,0,0,.55); color: #fff; padding: 6px 14px; border-radius: 14px; font-size: 13px; }}

.deck-grid {{ display: none; grid-template-columns: repeat(auto-fill, minmax(220px, 1fr)); gap: 16px; padding: 24px; }}
body.grid-mode .deck-grid {{ display: grid; }}
body.grid-mode .deck {{ display: none; }}
body.grid-mode .deck-nav {{ display: none; }}
.grid-tile {{
  aspect-ratio: 16/9; border: 1px solid #ccc; border-radius: 6px; overflow: hidden; position: relative;
  cursor: pointer; background: #fff;
}}
.grid-tile-inner {{ transform: scale(.2); transform-origin: top left; width: 500%; height: 500%; pointer-events: none; }}
.grid-tile-label {{
  position: absolute; bottom: 0; left: 0; right: 0; background: rgba(0,0,0,.65); color: #fff;
  font-size: 11px; padding: 3px 6px; white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
}}
"""


def _html_js() -> str:
    return """
(function () {
  var deck = document.getElementById('deck');
  var slides = Array.prototype.slice.call(deck.querySelectorAll('.slide'));
  var nav = document.getElementById('deck-nav');
  var grid = document.getElementById('deck-grid');
  var indicator = document.getElementById('deck-indicator');
  var current = 0;
  var presenting = false;
  var inGrid = false;

  function render() {
    slides.forEach(function (el, i) { el.classList.toggle('active', i === current); });
    indicator.textContent = (current + 1) + ' / ' + slides.length;
  }

  function enterPresenting(startAt) {
    presenting = true;
    if (typeof startAt === 'number') { current = startAt; }
    document.body.classList.add('presenting');
    render();
  }
  function exitPresenting() {
    presenting = false;
    document.body.classList.remove('presenting');
  }
  function goTo(i) {
    current = Math.max(0, Math.min(slides.length - 1, i));
    render();
  }
  function next() { if (current < slides.length - 1) { goTo(current + 1); } }
  function prev() { if (current > 0) { goTo(current - 1); } }

  function buildGrid() {
    grid.innerHTML = '';
    slides.forEach(function (slideEl, i) {
      var tile = document.createElement('div');
      tile.className = 'grid-tile';
      var inner = document.createElement('div');
      inner.className = 'grid-tile-inner';
      inner.innerHTML = slideEl.innerHTML;
      tile.appendChild(inner);
      var label = document.createElement('div');
      label.className = 'grid-tile-label';
      label.textContent = (i + 1) + '. ' + (slideEl.getAttribute('data-title') || '');
      tile.appendChild(label);
      tile.addEventListener('click', function () {
        exitGrid();
        enterPresenting(i);
      });
      grid.appendChild(tile);
    });
  }
  function enterGrid() {
    if (!grid.childElementCount) { buildGrid(); }
    inGrid = true;
    document.body.classList.add('grid-mode');
  }
  function exitGrid() {
    inGrid = false;
    document.body.classList.remove('grid-mode');
  }

  document.getElementById('btn-prev').addEventListener('click', prev);
  document.getElementById('btn-next').addEventListener('click', next);

  document.addEventListener('keydown', function (e) {
    if (e.key === 'Escape') {
      // These used to both fire on the same keypress: leaving presentation mode landed
      // straight in the grid overview instead of the plain deck view, because "not in
      // grid" was true the instant exitPresenting() cleared it. One Escape, one step back.
      if (presenting) { exitPresenting(); return; }
      if (inGrid) { exitGrid(); } else { enterGrid(); }
      return;
    }
    if (inGrid) { return; }
    if (!presenting) {
      if (e.key === 'ArrowRight' || e.key === 'ArrowLeft' || e.key === 'Enter') {
        enterPresenting(current);
      }
      return;
    }
    if (e.key === 'ArrowRight') { next(); }
    else if (e.key === 'ArrowLeft') { prev(); }
  });

  slides.forEach(function (el, i) {
    el.addEventListener('click', function () {
      if (!presenting && !inGrid) { enterPresenting(i); }
    });
  });

  var touchStartX = null;
  deck.addEventListener('touchstart', function (e) {
    touchStartX = e.changedTouches[0].clientX;
  }, { passive: true });
  deck.addEventListener('touchend', function (e) {
    if (touchStartX === null || !presenting) { return; }
    var dx = e.changedTouches[0].clientX - touchStartX;
    if (Math.abs(dx) > 40) { if (dx < 0) { next(); } else { prev(); } }
    touchStartX = null;
  }, { passive: true });

  render();
})();
"""


def deck_html(project: str) -> str:
    data = _require_deck(project)
    meta = data["meta"]
    theme = _theme_of(meta.get("theme", ""))
    slides = data["slides"]

    slides_html = "".join(_render_slide_html(s, meta, theme, i) for i, s in enumerate(slides))
    title = _esc(meta.get("title", "Deck"))

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>{_html_css(theme)}</style>
</head>
<body>
<div class="deck" id="deck">
{slides_html}
</div>
<div class="deck-nav" id="deck-nav">
  <button type="button" id="btn-prev" aria-label="Previous">&#8249;</button>
  <span id="deck-indicator">1 / {len(slides)}</span>
  <button type="button" id="btn-next" aria-label="Next">&#8250;</button>
</div>
<div class="deck-grid" id="deck-grid"></div>
<script>{_html_js()}</script>
</body>
</html>"""
