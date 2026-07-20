"""Text extraction from input documents (PDF, docx, xlsx, pptx, plain text).

Real client inputs are PDFs and Word files, not markdown: without extraction the agent
only sees a binary blob and can neither read nor cite it. Here the text is extracted
once at upload time and stored next to the file in `.extracted/`, so the skills always
read markdown and the model never has to know how to open a PDF.

Extracts are a cache, not a deliverable: they live in a hidden folder (excluded from
`/files`, from versioning and from backups) and can be deleted without losing anything.
"""
from __future__ import annotations

import re
from pathlib import Path

EXTRACTED_DIR = ".extracted"

# Extensions that are already text: read directly, no extract to generate.
TEXT_EXT = {".md", ".txt", ".csv", ".json", ".xml", ".yaml", ".yml"}

# Extensions we know how to extract text from. The rest (images, zips, executables)
# stays downloadable but not readable: that is correct, we do not try to guess.
BINARY_EXT = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}

SUPPORTED_EXT = TEXT_EXT | BINARY_EXT

# A 200-page brief is not needed whole: past this threshold the extract is truncated
# with an explicit notice, so the model's context does not blow up silently.
MAX_CHARS = 400_000


class MissingDependency(RuntimeError):
    """The library for that format is not installed."""


# --- extraction per format ---

def _from_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:  # pragma: no cover - depends on the environment
        raise MissingDependency("pypdf is not installed: pip install -r requirements.txt") from e
    reader = PdfReader(str(path))
    out: list[str] = []
    empty = True
    for n, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        if text:
            empty = False
        # The page marker is what allows citing "p. 7" and opening the PDF at
        # #page=7: without it the citation is not navigable.
        out.append(f"<!-- page {n} -->\n{text}")
    # A scanned PDF (pages that are images) would produce an extract made of markers
    # only: it would look readable while containing not one word. Better no extract at
    # all, so the UI flags it as unreadable and the user knows.
    return "" if empty else "\n\n".join(out)


def _from_docx(path: Path) -> str:
    try:
        import docx  # python-docx
    except ImportError as e:  # pragma: no cover
        raise MissingDependency("python-docx is not installed: pip install -r requirements.txt") from e
    doc = docx.Document(str(path))
    out: list[str] = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        # Word's "Heading N" styles become markdown headings: that is how the
        # document's chapters stay citable as anchors.
        # "Titolo N" is the same style under an Italian Word install: matching both
        # means a document written on a localised Word still yields its chapters.
        style = (p.style.name or "") if p.style is not None else ""
        m = re.match(r"Heading (\d)", style) or re.match(r"Titolo (\d)", style)
        if m:
            out.append(f"{'#' * min(int(m.group(1)), 6)} {text}")
        else:
            out.append(text)
    for table in doc.tables:
        out.append(_table_md([[c.text.strip() for c in row.cells] for row in table.rows]))
    return "\n\n".join(x for x in out if x)


def _from_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    # read_only mode keeps only a window of rows in memory instead of the whole sheet,
    # which is exactly why it also needs an explicit close() to release the file — and
    # that close() used to sit only on the happy path. extract_text() swallows every
    # exception from this function, so a corrupt workbook that raised anywhere between
    # here and the end used to leave the handle open for the rest of the process — on
    # Windows, long enough that the document could no longer be deleted or replaced.
    wb = load_workbook(str(path), data_only=True, read_only=True)
    try:
        out: list[str] = []
        for ws in wb.worksheets:
            rows = [[("" if c is None else str(c)) for c in row]
                    for row in ws.iter_rows(values_only=True)]
            rows = [r for r in rows if any(x.strip() for x in r)]
            if not rows:
                continue
            out.append(f"## Sheet: {ws.title}\n\n{_table_md(rows)}")
        return "\n\n".join(out)
    finally:
        wb.close()


def _from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover
        raise MissingDependency("python-pptx is not installed: pip install -r requirements.txt") from e
    prs = Presentation(str(path))
    out: list[str] = []
    for n, slide in enumerate(prs.slides, start=1):
        parts = [sh.text.strip() for sh in slide.shapes
                 if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        out.append(f"## Slide {n}\n\n" + "\n\n".join(parts))
    return "\n\n".join(out)


def _table_md(rows: list[list[str]]) -> str:
    """Markdown table. Cells containing a pipe escape it, otherwise the table breaks."""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    def line(cells: list[str]) -> str:
        cells = [c.replace("|", "\\|").replace("\n", " ") for c in cells]
        cells += [""] * (width - len(cells))
        return "| " + " | ".join(cells) + " |"
    head, *body = rows
    return "\n".join([line(head), "| " + " | ".join(["---"] * width) + " |",
                      *(line(r) for r in body)])


_EXTRACTORS = {".pdf": _from_pdf, ".docx": _from_docx, ".xlsx": _from_xlsx,
               ".xlsm": _from_xlsx, ".pptx": _from_pptx}


def extract_text(path: Path) -> str | None:
    """The document's text, or None if the format is not supported.

    Never raises on a corrupt file: an unreadable document must not fail the upload,
    it simply ends up without an extract.
    """
    ext = path.suffix.lower()
    try:
        if ext in TEXT_EXT:
            return path.read_text(encoding="utf-8", errors="replace")
        fn = _EXTRACTORS.get(ext)
        if fn is None:
            return None
        text = fn(path)
    except MissingDependency:
        raise
    except Exception:
        return None
    text = (text or "").strip()
    if not text:
        return None
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + "\n\n<!-- document truncated: past the reading limit -->"
    return text


# --- extracts on disk ---

def extract_path(project_dir: Path, rel: str) -> Path:
    """`docs/flows.pdf` -> `.extracted/docs__flows.pdf.md`. Flat path: the folder stays
    inspectable by eye and no parallel trees get created."""
    return project_dir / EXTRACTED_DIR / (rel.replace("/", "__").replace("\\", "__") + ".md")


def needs_extract(rel: str) -> bool:
    return Path(rel).suffix.lower() in BINARY_EXT


def build_extract(project_dir: Path, rel: str) -> Path | None:
    """Generates (or regenerates) the extract of a project file. None if not applicable."""
    src = project_dir / rel
    if not src.is_file() or not needs_extract(rel):
        return None
    try:
        text = extract_text(src)
    except MissingDependency:
        return None
    if text is None:
        return None
    dest = extract_path(project_dir, rel)
    dest.parent.mkdir(parents=True, exist_ok=True)
    header = f"<!-- automatic extract of {rel} — do not edit by hand -->\n\n"
    dest.write_text(header + text, encoding="utf-8")
    return dest


def drop_extract(project_dir: Path, rel: str) -> None:
    extract_path(project_dir, rel).unlink(missing_ok=True)


def readable_text(project_dir: Path, rel: str) -> str | None:
    """Text of a project file: directly if it already is text, from the extract otherwise."""
    src = project_dir / rel
    if not src.is_file():
        return None
    if src.suffix.lower() in TEXT_EXT:
        return src.read_text(encoding="utf-8", errors="replace")
    ext = extract_path(project_dir, rel)
    if not ext.is_file():
        build_extract(project_dir, rel)
    if ext.is_file():
        return ext.read_text(encoding="utf-8", errors="replace")
    return None


# --- brief: can be md, pdf or docx ---

# Order of preference: if for some reason there were two brief.* files, markdown wins
# (it is the one we write ourselves when the project starts from discovery).
BRIEF_EXT_ORDER = [".md", ".docx", ".pdf", ".txt"]


def brief_file(project_dir: Path) -> str | None:
    """Name of the brief file present ('brief.pdf', 'brief.md', ...) or None."""
    for ext in BRIEF_EXT_ORDER:
        if (project_dir / f"brief{ext}").is_file():
            return f"brief{ext}"
    return None


def brief_text(project_dir: Path) -> str | None:
    rel = brief_file(project_dir)
    return readable_text(project_dir, rel) if rel else None


# --- anchors for citations ---

_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*#*$")
_PAGE_RE = re.compile(r"^<!--\s*page\s+(\d+)\s*-->$")


def slugify(text: str) -> str:
    s = re.sub(r"[^\w\s-]", "", text.lower(), flags=re.UNICODE)
    return re.sub(r"[\s_]+", "-", s).strip("-")


def headings(text: str) -> list[dict]:
    """Citable chapters of a markdown text, with slug, line and source page.

    The page is what PDFs need: the navigable anchor there is `#page=N`, not the slug.
    """
    out: list[dict] = []
    page = 0
    seen: dict[str, int] = {}
    for n, line in enumerate(text.splitlines(), start=1):
        pm = _PAGE_RE.match(line.strip())
        if pm:
            page = int(pm.group(1))
            continue
        m = _HEADING_RE.match(line.rstrip())
        if not m:
            continue
        title = m.group(2).strip()
        slug = slugify(title) or f"section-{n}"
        # Two chapters with the same title must get different slugs, otherwise the
        # citation always lands on the first one.
        if slug in seen:
            seen[slug] += 1
            slug = f"{slug}-{seen[slug]}"
        else:
            seen[slug] = 0
        out.append({"level": len(m.group(1)), "title": title, "slug": slug,
                    "line": n, "page": page or None})
    return out
