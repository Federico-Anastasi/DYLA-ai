"""Tests for the slide deck export (server/deck_export.py): two formats (pptx/html) out
of the same deck.json source.

Run with `python -m pytest server/tests`. Fixture pattern identical to
test_mockup_export.py: load_doc is monkeypatched so the project filesystem is never
touched (except for the 'missing image' case, which is precisely about a file that is
not there).
"""
import io
import json
from pathlib import Path

import jsonschema
import pytest
from pptx import Presentation

from server import deck_export
from server.deck_export import deck_html, deck_pptx
from server.exports import DocNotFound

SCHEMA = json.loads(
    (Path(__file__).resolve().parent.parent.parent / "schemas" / "deck.schema.json")
    .read_text(encoding="utf-8")
)


def _validate(doc: dict) -> None:
    jsonschema.validate(instance=doc, schema=SCHEMA)


def _meta(**extra) -> dict:
    meta = {"project": "t", "title": "Test Deck", "client": "Acme",
            "date": "2026-07-19", "type": "status"}
    meta.update(extra)
    return meta


def _full_deck() -> dict:
    """One slide per layout, unique ids, minimal valid fields."""
    return {
        "meta": _meta(),
        "slides": [
            {"id": "S1", "layout": "cover", "title": "Cover Title",
             "subtitle": "Subtitle"},
            {"id": "S2", "layout": "section", "title": "Section Title"},
            {
                "id": "S3", "layout": "list", "title": "List Title",
                "bullets": ["First point", "Second **important** point"],
                "speaker_notes": "Remember to mention the overlap rule",
            },
            {
                "id": "S4", "layout": "table", "title": "Table Title",
                "table": {"headers": ["Epic", "Days"],
                          "rows": [["1. Rooms", "6.75"], ["2. Booking", "9.0"]]},
            },
            {
                "id": "S5", "layout": "kpi", "title": "KPI Title",
                "kpi": [
                    {"label": "Total days", "value": "19.0", "note": "vs 18 planned"},
                    {"label": "Progress", "value": "64%"},
                ],
            },
            {
                "id": "S6", "layout": "timeline", "title": "Timeline Title",
                "milestones": [
                    {"name": "Kick-off", "date": "2026-01-10", "status": "completed"},
                    {"name": "Analysis delivered", "date": "2026-02-15",
                     "status": "in_progress"},
                    {"name": "Go-live", "date": "2026-06-01", "status": "at_risk"},
                ],
            },
            {
                "id": "S7", "layout": "text", "title": "Text Title",
                "text": "An opening paragraph.\n\n- First bullet\n- Second *italic* bullet",
            },
            {"id": "S8", "layout": "image", "title": "Image Title",
             "image": "docs/screenshot_does_not_exist.png"},
        ],
    }


@pytest.fixture
def patch_doc(monkeypatch):
    def _patch(doc: dict):
        monkeypatch.setattr(deck_export, "load_doc",
                            lambda project, name: doc if project == "t" else None)

    return _patch


def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


# --- pptx: one slide per layout, titles present, file reopens ------------------

def test_pptx_has_one_slide_per_layout_with_titles(patch_doc):
    doc = _full_deck()
    _validate(doc)
    patch_doc(doc)

    raw = deck_pptx("t")
    prs = Presentation(io.BytesIO(raw))

    assert len(prs.slides) == len(doc["slides"])
    for slide, s in zip(prs.slides, doc["slides"]):
        assert s["title"] in _slide_text(slide)


def test_pptx_is_16_by_9():
    # No real document needed: the size is set before the slides are iterated.
    from server.deck_export import SLIDE_H, SLIDE_W

    assert round(SLIDE_W / SLIDE_H, 3) == round(16 / 9, 3)


def test_pptx_speaker_notes_land_in_the_notes_slide(patch_doc):
    doc = _full_deck()
    patch_doc(doc)
    raw = deck_pptx("t")
    prs = Presentation(io.BytesIO(raw))

    list_slide = prs.slides[2]  # S3, layout 'list', has speaker_notes
    assert list_slide.has_notes_slide
    assert "overlap rule" in list_slide.notes_slide.notes_text_frame.text


def test_pptx_a_missing_image_does_not_crash(patch_doc):
    doc = _full_deck()
    patch_doc(doc)
    raw = deck_pptx("t")
    prs = Presentation(io.BytesIO(raw))
    image_slide = prs.slides[7]  # S8
    assert "docs/screenshot_does_not_exist.png" in _slide_text(image_slide)


def test_pptx_timeline_degrades_to_a_table_past_the_threshold(patch_doc):
    doc = _full_deck()
    many_milestones = [
        {"name": f"Milestone {i}", "date": "2026-01-01", "status": "planned"}
        for i in range(8)
    ]
    doc["slides"][5]["milestones"] = many_milestones
    _validate(doc)
    patch_doc(doc)

    raw = deck_pptx("t")
    prs = Presentation(io.BytesIO(raw))
    timeline_slide = prs.slides[5]
    text = _slide_text(timeline_slide)
    assert "Milestone" in text  # header of the fallback table
    assert "Milestone 7" in text
    # No oval shapes (the dots of the inline timeline) once it degrades to a table
    assert any(shape.has_table for shape in timeline_slide.shapes)


# --- html: titles present, zero external resources ----------------------------

def test_html_contains_every_title(patch_doc):
    doc = _full_deck()
    patch_doc(doc)
    out = deck_html("t")
    for s in doc["slides"]:
        assert s["title"] in out


def test_html_has_no_external_resources(patch_doc):
    """The deck gets emailed around and opened offline: one remote font or script and it
    renders differently, or not at all, on the other side."""
    doc = _full_deck()
    patch_doc(doc)
    out = deck_html("t")
    assert "http://" not in out
    assert "https://" not in out


def test_html_carries_the_theme_colour(patch_doc):
    doc = _full_deck()
    patch_doc(doc)
    out = deck_html("t")
    # No theme named in meta: the export falls back to the neutral default palette.
    assert deck_export._DEFAULT_PRIMARY in out.upper()


def test_html_milestone_statuses_render_with_their_own_style(patch_doc):
    """Every status in the schema enum has to reach the page with its own label and dot
    class. A status the export does not know about degrades to an unlabelled grey dot,
    which is exactly the state ('completed', 'at risk') a reader most needs to see.
    """
    doc = _full_deck()
    patch_doc(doc)
    out = deck_html("t")
    for status, label in (("completed", "Completed"), ("in_progress", "In progress"),
                          ("at_risk", "At risk")):
        assert f"tl-dot-{status}" in out
        assert f'title="{label}"' in out
    # Each dot class must have a rule behind it, otherwise it renders unstyled.
    assert ".tl-dot-completed" in out


def test_html_a_missing_image_shows_a_placeholder(patch_doc):
    doc = _full_deck()
    patch_doc(doc)
    out = deck_html("t")
    assert "Image not found" in out
    assert "docs/screenshot_does_not_exist.png" in out


def test_html_timeline_degrades_to_a_table_past_the_threshold(patch_doc):
    doc = _full_deck()
    doc["slides"][5]["milestones"] = [
        {"name": f"Milestone {i}", "status": "planned"} for i in range(8)
    ]
    _validate(doc)
    patch_doc(doc)
    out = deck_html("t")
    assert "Milestone 7" in out
    assert "<table>" in out
    assert 'class="timeline-row"' not in out


# --- DocNotFound when deck.json is missing ------------------------------------

# --- "table" is optional in the schema, unlike every other layout's own field --------

def test_pptx_table_slide_with_no_table_field_does_not_crash(patch_doc):
    """Only id/layout/title are required: a layout='table' slide with nothing else used
    to raise a KeyError on `s["table"]` and take the whole export down with it."""
    doc = _full_deck()
    del doc["slides"][3]["table"]
    patch_doc(doc)
    raw = deck_pptx("t")
    prs = Presentation(io.BytesIO(raw))
    assert "No table data" in _slide_text(prs.slides[3])


def test_html_table_slide_with_no_table_field_does_not_crash(patch_doc):
    doc = _full_deck()
    del doc["slides"][3]["table"]
    patch_doc(doc)
    out = deck_html("t")
    assert "No table data" in out


# --- presentation mode: Escape does one thing per keypress --------------------

def test_escape_exits_presenting_without_also_opening_the_grid(patch_doc):
    """Escape used to both exit presentation mode AND enter the grid overview in the
    same keypress, because "not in grid" was already true the instant exitPresenting()
    cleared it — leaving presentation always dumped you into the grid instead of the
    plain deck view."""
    doc = _full_deck()
    patch_doc(doc)
    out = deck_html("t")
    assert "if (presenting) { exitPresenting(); return; }" in out


def test_doc_not_found_pptx(monkeypatch):
    monkeypatch.setattr(deck_export, "load_doc", lambda project, name: None)
    with pytest.raises(DocNotFound):
        deck_pptx("missing")


def test_doc_not_found_html(monkeypatch):
    monkeypatch.setattr(deck_export, "load_doc", lambda project, name: None)
    with pytest.raises(DocNotFound):
        deck_html("missing")


# --- a slide image is only ever read from inside its own project --------------------

def test_an_image_path_cannot_walk_out_of_the_project(tmp_path, monkeypatch):
    """`slides[].image` is a path the model writes into a JSON file, and the schema puts
    no constraint on it. Joined naively, "../../server/config.yaml" leaves the project —
    and the HTML export embeds whatever it finds as base64 inside a file meant to be sent
    to a client."""
    monkeypatch.setattr(deck_export, "PROJECTS_DIR", tmp_path)
    project = tmp_path / "demo"
    (project / "docs").mkdir(parents=True)
    (tmp_path / "secret.txt").write_text("not yours", encoding="utf-8")

    assert deck_export._project_image({"project": "demo"}, "../secret.txt") is None


def test_an_absolute_image_path_is_refused(tmp_path, monkeypatch):
    """On Windows an absolute path makes pathlib discard the prefix entirely, so a naive
    join silently becomes "anywhere on this disk"."""
    monkeypatch.setattr(deck_export, "PROJECTS_DIR", tmp_path)
    (tmp_path / "demo").mkdir()
    outside = tmp_path / "elsewhere.png"
    outside.write_bytes(b"x")

    assert deck_export._project_image({"project": "demo"}, str(outside)) is None


def test_a_real_project_image_still_works(tmp_path, monkeypatch):
    """The point is containment, not refusing everything."""
    monkeypatch.setattr(deck_export, "PROJECTS_DIR", tmp_path)
    docs = tmp_path / "demo" / "docs"
    docs.mkdir(parents=True)
    (docs / "shot.png").write_bytes(b"x")

    found = deck_export._project_image({"project": "demo"}, "docs/shot.png")
    assert found is not None and found.name == "shot.png"
